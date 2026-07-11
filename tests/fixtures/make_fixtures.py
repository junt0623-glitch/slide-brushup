#!/usr/bin/env python3
"""
テスト用pptxフィクスチャ生成スクリプト（python-pptx使用）。
GitHub Actions上でNode.jsのpptx生成ライブラリに頼らず済むよう、
このスクリプトの「実行結果」であるバイナリ.pptxファイルをリポジトリに
そのままコミットする運用とする（CI側での再生成は不要）。

生成物:
  - basic.pptx     : bt01（テキスト抽出）/ bt03（rels経由のスライド順解決）用
  - image.pptx      : bt02（画像バイトハッシュ一致）用
  - preserve.pptx   : bt04（温存フラグ検出）用
  - expectations.json : 上記フィクスチャの「正解」をPlaywrightテストが参照する
"""
import hashlib
import json
import os
import re
import zipfile

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu, Inches, Pt

OUT = os.path.dirname(os.path.abspath(__file__))

SLIDE_W = 12192000  # 16:9, EMU
SLIDE_H = 6858000


def _parse_rels(rels_xml: str) -> dict:
    out = {}
    for tag in re.findall(r"<Relationship\b[^>]*/>", rels_xml):
        idm = re.search(r'Id="([^"]+)"', tag)
        tgtm = re.search(r'Target="([^"]+)"', tag)
        if idm and tgtm:
            out[idm.group(1)] = tgtm.group(1)
    return out


def _get_slide_order_filenames(path: str):
    """presentation.xml + presentation.xml.rels から、rels経由での論理スライド順を求める。
    JS側パーサーが行うのと同じ解決手順をPythonでも踏むことで、
    フィクスチャ自体の正しさを保証する。"""
    with zipfile.ZipFile(path) as z:
        pres_xml = z.read("ppt/presentation.xml").decode("utf-8")
        rels_xml = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
    rid_to_target = _parse_rels(rels_xml)
    m = re.search(r"<p:sldIdLst>(.*?)</p:sldIdLst>", pres_xml, re.S)
    rids = re.findall(r'r:id="([^"]+)"', m.group(1))
    return [rid_to_target[r] for r in rids]


def _rewrite_zip(path: str, mutate):
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        data = {n: z.read(n) for n in names}
    mutate(data)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for n in names:
            z.writestr(n, data[n])


def _reverse_first_last_slide_order(path: str):
    """sldIdLst内の r:id を並べ替え、『ファイル名の物理順』と『論理表示順』を
    意図的にズラす（bt03: rels経由の正しい解決を検証するため）。"""

    def mutate(data):
        xml_text = data["ppt/presentation.xml"].decode("utf-8")
        m = re.search(r"<p:sldIdLst>(.*?)</p:sldIdLst>", xml_text, re.S)
        block = m.group(1)
        rids = re.findall(r'r:id="([^"]+)"', block)
        assert len(rids) == 3, rids
        new_rids = [rids[2], rids[1], rids[0]]
        tags = re.findall(r"<p:sldId[^/]*/>", block)
        assert len(tags) == 3
        new_tags = [
            re.sub(r'r:id="[^"]+"', f'r:id="{nr}"', tag)
            for tag, nr in zip(tags, new_rids)
        ]
        new_block = "".join(new_tags)
        new_xml = xml_text[: m.start(1)] + new_block + xml_text[m.end(1) :]
        data["ppt/presentation.xml"] = new_xml.encode("utf-8")

    _rewrite_zip(path, mutate)


def _inject_fake_diagram(path: str, logical_index: int):
    """指定した論理スライド位置(0始まり)のXMLに、SmartArt(diagram)を模した
    graphicFrameを直接注入する（python-pptxはSmartArt生成に非対応のため）。
    参照先rIdは実在しないダミーだが、本アプリの検出ロジックは
    graphicData uri文字列の走査のみで判定するため実害はない。"""
    order = _get_slide_order_filenames(path)
    target_file = "ppt/" + order[logical_index]

    def mutate(data):
        xml_text = data[target_file].decode("utf-8")
        stub = (
            "<p:graphicFrame>"
            '<p:nvGraphicFramePr><p:cNvPr id="999" name="TestDiagram"/>'
            "<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>"
            '<p:xfrm><a:off x="1000000" y="1000000"/><a:ext cx="3000000" cy="3000000"/></p:xfrm>'
            "<a:graphic><a:graphicData "
            'uri="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
            '<dgm:relIds xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" '
            'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
            'r:dm="rIdFake1" r:lo="rIdFake2" r:qs="rIdFake3" r:cs="rIdFake4"/>'
            "</a:graphicData></a:graphic></p:graphicFrame>"
        )
        assert "</p:spTree>" in xml_text
        xml_text = xml_text.replace("</p:spTree>", stub + "</p:spTree>")
        data[target_file] = xml_text.encode("utf-8")

    _rewrite_zip(path, mutate)


def _new_prs():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    return prs


def build_basic():
    prs = _new_prs()
    layout = prs.slide_layouts[6]  # 白紙
    contents = [
        ("フィクスチャA", "これはスライドAの本文です。改行テスト。\n二行目のテキスト。"),
        ("フィクスチャB", "スライドBの本文、太字ランを含みます。"),
        ("フィクスチャC", "スライドCの本文、これが最後のスライドです。"),
    ]
    for title, body in contents:
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        r.font.size = Pt(32)
        r.font.bold = True
        bb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
        bb.text_frame.text = body
    path = os.path.join(OUT, "basic.pptx")
    prs.save(path)
    _reverse_first_last_slide_order(path)
    return path


def build_image():
    prs = _new_prs()
    layout = prs.slide_layouts[6]
    img_path = os.path.join(OUT, "_tmp_pixel.png")
    Image.new("RGB", (60, 40), (91, 152, 133)).save(img_path, "PNG")
    with open(img_path, "rb") as f:
        img_bytes = f.read()
    img_hash = hashlib.sha256(img_bytes).hexdigest()

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(6), Inches(0.8)
    ).text_frame.text = "画像フィクスチャ"
    slide1.shapes.add_picture(
        img_path, Inches(1), Inches(1.5), width=Inches(2), height=Inches(1.333)
    )

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(6), Inches(0.8)
    ).text_frame.text = "画像フィクスチャ（同一画像の再利用）"
    slide2.shapes.add_picture(
        img_path, Inches(3), Inches(2), width=Inches(2), height=Inches(1.333)
    )

    path = os.path.join(OUT, "image.pptx")
    prs.save(path)
    os.remove(img_path)
    return path, img_hash


def build_preserve():
    prs = _new_prs()
    layout = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(layout)
    s1.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(6), Inches(0.6)
    ).text_frame.text = "グラフを含むスライド"
    cd = CategoryChartData()
    cd.categories = ["一月", "二月", "三月"]
    cd.add_series("売上", (10, 20, 30))
    s1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.2), Inches(6), Inches(4), cd
    )

    s2 = prs.slides.add_slide(layout)
    s2.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(6), Inches(0.6)
    ).text_frame.text = "表を含むスライド"
    tbl = s2.shapes.add_table(2, 2, Inches(1), Inches(1.2), Inches(6), Inches(1.5)).table
    tbl.cell(0, 0).text = "項目"
    tbl.cell(0, 1).text = "値"
    tbl.cell(1, 0).text = "テスト行"
    tbl.cell(1, 1).text = "123"

    s3 = prs.slides.add_slide(layout)
    s3.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(6), Inches(0.6)
    ).text_frame.text = "通常のテキストのみ（対照群）"

    s4 = prs.slides.add_slide(layout)
    s4.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(6), Inches(0.6)
    ).text_frame.text = "SmartArt想定スライド"

    path = os.path.join(OUT, "preserve.pptx")
    prs.save(path)
    _inject_fake_diagram(path, logical_index=3)
    return path


def build_shapes():
    """図形ツールの認識を検証するための構成。
    4つ目の図形は塗りを明示指定しない＝PowerPointの図形ツールで作ったままの状態
    （<p:style>のテーマ参照のみ）。これがユーザー報告の「図形が全て消える」不具合の
    再現ケースであり、テーマ色解決の回帰テストを兼ねる。"""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    prs = _new_prs()
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    r1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    r1.fill.solid()
    r1.fill.fore_color.rgb = RGBColor(0x1A, 0x7A, 0x5E)
    r1.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    r1.line.width = Pt(2)

    a1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4), Inches(1), Inches(1.5), Inches(0.8))
    a1.fill.solid()
    a1.fill.fore_color.rgb = RGBColor(0xC0, 0x8A, 0x3E)

    r2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(3), Inches(1))
    r2.fill.solid()
    r2.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    r2.text_frame.text = "ラベル付き図形"

    # 塗り未指定 = テーマスタイル参照のみ（accent1 = 4F81BD が期待値）
    slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(3), Inches(1.5), Inches(1.5))

    path = os.path.join(OUT, "shapes.pptx")
    prs.save(path)
    return path


ANIM_TIMING_STUB = (
    "<p:timing><p:tnLst><p:par>"
    '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
    "<p:childTnLst><p:seq concurrent=\"1\" nextAc=\"seek\">"
    '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
    '<p:par><p:cTn id="3" fill="hold"><p:childTnLst>'
    '<p:par><p:cTn id="4" fill="hold"><p:childTnLst>'
    '<p:par><p:cTn id="5" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect">'
    "<p:childTnLst>"
    '<p:set><p:cBhvr><p:cTn id="6" dur="1" fill="hold"/>'
    '<p:tgtEl><p:spTgt spid="2"/></p:tgtEl>'
    "<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>"
    "<p:to><p:strVal val=\"visible\"/></p:to></p:set>"
    "</p:childTnLst></p:cTn></p:par>"
    "</p:childTnLst></p:cTn></p:par>"
    "</p:childTnLst></p:cTn></p:par>"
    "</p:childTnLst></p:cTn></p:seq></p:childTnLst>"
    "</p:cTn></p:par></p:tnLst></p:timing>"
)


def build_anim():
    """アニメーション・画面切り替えの検出/引き継ぎ検証用。
    python-pptxはアニメーション生成に非対応のため、実際のPowerPointが出力する
    形式に沿ったXML断片をスライドへ直接注入する。
      slide1: アニメーションあり（温存の初期値=ONになるはず）
      slide2: 画面切り替えのみ（再構築されても切り替えは引き継がれるはず）
      slide3: どちらもなし（対照群）
    """
    prs = _new_prs()
    layout = prs.slide_layouts[6]
    for label in ("アニメーション付きスライド", "画面切り替え付きスライド", "効果なしスライド"):
        s = prs.slides.add_slide(layout)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        tb.text_frame.text = label
    path = os.path.join(OUT, "anim.pptx")
    prs.save(path)

    def mutate(data):
        # slide1: timing(アニメーション)を注入
        xml1 = data["ppt/slides/slide1.xml"].decode("utf-8")
        assert "</p:sld>" in xml1
        data["ppt/slides/slide1.xml"] = xml1.replace(
            "</p:sld>", ANIM_TIMING_STUB + "</p:sld>"
        ).encode("utf-8")
        # slide2: transition(画面切り替え・フェード)を注入
        xml2 = data["ppt/slides/slide2.xml"].decode("utf-8")
        data["ppt/slides/slide2.xml"] = xml2.replace(
            "</p:sld>", '<p:transition spd="slow"><p:fade/></p:transition></p:sld>'
        ).encode("utf-8")

    _rewrite_zip(path, mutate)

    # 注入後のXMLが妥当であることを確認（DOMParser相当）
    import xml.etree.ElementTree as ET
    with zipfile.ZipFile(path) as z:
        for n in ("ppt/slides/slide1.xml", "ppt/slides/slide2.xml"):
            ET.fromstring(z.read(n).decode("utf-8"))
    return path


def main():
    build_basic()
    _img_path, img_hash = build_image()
    build_preserve()
    build_shapes()
    build_anim()

    expectations = {
        "basic.pptx": {
            "slideCount": 3,
            "titlesInOrder": ["フィクスチャC", "フィクスチャB", "フィクスチャA"],
            "note": "ファイル名(slide1/2/3.xml)は元のA,B,C順のままだが、"
            "sldIdLstの並びをC,B,Aに変更してある。rels経由で解決すればC,B,A、"
            "ファイル名ソートで誤解決するとA,B,Cになる。",
        },
        "image.pptx": {
            "slideCount": 2,
            "imageSha256": img_hash,
            "note": "2枚のスライドが同一画像を参照する（重複排除の確認にも使える）。",
        },
        "preserve.pptx": {
            "slideCount": 4,
            "expectedPreserve": [
                {"index": 0, "preserve": True, "reason": "chart"},
                {"index": 1, "preserve": False, "reason": None, "note": "表は再構築対象のため温存不要"},
                {"index": 2, "preserve": False, "reason": None},
                {"index": 3, "preserve": True, "reason": "smartart"},
            ],
        },
    }
    with open(os.path.join(OUT, "expectations.json"), "w", encoding="utf-8") as f:
        json.dump(expectations, f, ensure_ascii=False, indent=2)

    for name in ("basic.pptx", "image.pptx", "preserve.pptx", "shapes.pptx", "anim.pptx", "expectations.json"):
        p = os.path.join(OUT, name)
        print(f"{name}: {os.path.getsize(p)} bytes")


if __name__ == "__main__":
    main()
